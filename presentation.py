"""
AgentGate — ETHGlobal Cannes 2026 Presentation
Generates a dark-themed PowerPoint with World (purple) + Hedera (green) tags
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
BG = RGBColor(0x0A, 0x0A, 0x0A)
WHITE = RGBColor(0xE5, 0xE7, 0xEB)
GRAY = RGBColor(0x88, 0x88, 0x88)
DARK_GRAY = RGBColor(0x44, 0x44, 0x44)
GREEN = RGBColor(0x4A, 0xDE, 0x80)
PURPLE = RGBColor(0x82, 0x59, 0xEF)
HEDERA_GREEN = RGBColor(0x00, 0xD4, 0x6E)
ACCENT = RGBColor(0x8B, 0x5C, 0xF6)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG

def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_tag(slide, left, top, text, color):
    """Small colored badge"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(1.1), Inches(0.32))
    shape.fill.solid()
    rgb_int = int.from_bytes(bytes(color), 'big')
    r, g, b = (rgb_int >> 16) & 0xFF, (rgb_int >> 8) & 0xFF, rgb_int & 0xFF
    shape.fill.fore_color.rgb = RGBColor(r // 5, g // 5, b // 5)
    shape.line.color.rgb = color
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

def add_bullet_slide(slide, left, top, width, bullets, size=16, color=GRAY):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(8)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)
add_text(slide, 1.5, 1.5, 10, 1.2, "AgentGate", size=54, bold=True, color=WHITE)
add_text(slide, 1.5, 2.7, 10, 0.8, "Pay-per-call API Marketplace for Human-Backed AI Agents", size=22, color=ACCENT)
add_text(slide, 1.5, 3.8, 10, 0.5, "ETHGlobal Cannes 2026", size=16, color=DARK_GRAY)
add_text(slide, 1.5, 4.5, 10, 0.5, "Built on Hedera  ×  World AgentKit  ×  x402 Protocol", size=14, color=GRAY)
add_tag(slide, 1.5, 5.5, "🟣 World", PURPLE)
add_tag(slide, 2.8, 5.5, "🟢 Hedera", HEDERA_GREEN)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 1.5, 0.8, 10, 0.8, "The Problem", size=36, bold=True, color=WHITE)
add_bullet_slide(slide, 1.5, 2.0, 10, [
    "AI agents need APIs (OpenAI, data feeds, tools) but can't subscribe or pay like humans",
    "API publishers have no way to monetize per-call without building payment infrastructure",
    "No way to distinguish human-backed agents from spam bots at the API level",
    "Without identity: rate abuse, no trust, no fair pricing",
    "Result: walled gardens, no open agent economy",
], size=18, color=GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — The Solution
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 1.5, 0.8, 10, 0.8, "AgentGate — The Solution", size=36, bold=True, color=WHITE)
add_text(slide, 1.5, 2.0, 10, 0.6, "Any publisher lists their API. Any agent pays per call. Human identity optional.", size=18, color=ACCENT)

# Left column — Publisher
add_text(slide, 1.5, 3.0, 4.5, 0.5, "Publisher", size=22, bold=True, color=GREEN)
add_bullet_slide(slide, 1.5, 3.6, 4.5, [
    "→ Provides backend URL + API key",
    "→ Sets price per call (USD → HBAR)",
    "→ Optionally requires WorldID",
    "→ Gets an AgentGate proxy URL",
    "→ API key never exposed to agents",
], size=14, color=GRAY)

# Right column — Agent
add_text(slide, 7, 3.0, 5, 0.5, "AI Agent", size=22, bold=True, color=PURPLE)
add_bullet_slide(slide, 7, 3.6, 5, [
    "→ Discovers endpoint in registry",
    "→ Calls proxy URL → gets 402 challenge",
    "→ Pays HBAR on Hedera (3s finality)",
    "→ If WorldID verified: 3 free calls first",
    "→ Gets API response — zero setup needed",
], size=14, color=GRAY)

add_tag(slide, 1.5, 6.5, "🟣 World", PURPLE)
add_tag(slide, 2.8, 6.5, "🟢 Hedera", HEDERA_GREEN)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — World AgentKit (for World Prize)
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_tag(slide, 11.5, 0.5, "🟣 World", PURPLE)
add_text(slide, 1.5, 0.8, 10, 0.8, "World AgentKit Integration", size=36, bold=True, color=PURPLE)
add_text(slide, 1.5, 1.8, 10, 0.5, "Human-backed agents get privileged access", size=18, color=GRAY)

add_bullet_slide(slide, 1.5, 2.8, 5.5, [
    "Per-endpoint WorldID toggle — publishers decide",
    "createAgentBookVerifier() → real World Chain lookup",
    "SIWE signing via formatSIWEMessage()",
    "Verified agents → 3 free API calls (free-trial)",
    "After free-trial → pay HBAR like everyone else",
    "Unverified agents → pay from call #1 or blocked",
    "InMemoryAgentKitStorage pattern (DB-ready)",
], size=15, color=GRAY)

# Right side — flow
add_text(slide, 7.5, 2.8, 5, 0.4, "Agent Flow:", size=14, bold=True, color=PURPLE)
add_bullet_slide(slide, 7.5, 3.3, 5, [
    "1. Agent sends agentkit header (SIWE proof)",
    "2. Server validates signature",
    "3. AgentBook lookup on World Chain",
    "4. ✅ In AgentBook → free-trial (3 calls)",
    "5. ❌ Not in AgentBook → must pay HBAR",
    "6. Publisher controls: WorldID on/off per endpoint",
], size=13, color=GRAY)

add_text(slide, 7.5, 5.8, 5, 0.4, "Not mocked — real AgentBook verification on World Chain", size=12, bold=True, color=GREEN)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — x402 on Hedera (for Hedera Prize)
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_tag(slide, 11.5, 0.5, "🟢 Hedera", HEDERA_GREEN)
add_text(slide, 1.5, 0.8, 10, 0.8, "x402 Payments on Hedera", size=36, bold=True, color=HEDERA_GREEN)
add_text(slide, 1.5, 1.8, 10, 0.5, "Real HBAR micropayments — not simulated", size=18, color=GRAY)

add_bullet_slide(slide, 1.5, 2.8, 5.5, [
    "HTTP 402 Payment Required — industry standard",
    "Prices in USD, settled in HBAR (live Mirror Node rate)",
    "Payment verification via Mirror Node REST API",
    "mirror.hedera.com/api/v1/contracts/results/{txHash}",
    "3-second finality — agent pays and gets response in ~5s",
    "Replay protection (used tx hash tracking)",
    "All transactions verifiable on HashScan",
], size=15, color=GRAY)

# Right side — numbers
add_text(slide, 7.5, 2.8, 5, 0.4, "Live on Hedera Testnet:", size=14, bold=True, color=HEDERA_GREEN)
add_bullet_slide(slide, 7.5, 3.3, 5, [
    "PublisherRegistry: 0xFBCee3E3...89",
    "AgentGatePaymaster: 0xfbC79b8d...a1",
    "EntryPoint v0.7: 0x00000000717...32",
    "",
    "Real HBAR spent in demo:",
    "  Weather API: ~0.115 HBAR/call ($0.01)",
    "  Price Feed:  ~0.058 HBAR/call ($0.005)",
    "  OpenAI proxy: ~2.3 HBAR/call ($0.20)",
], size=13, color=GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Smart Contracts
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_tag(slide, 11.5, 0.5, "🟢 Hedera", HEDERA_GREEN)
add_text(slide, 1.5, 0.8, 10, 0.8, "Smart Contracts", size=36, bold=True, color=WHITE)

# Left — Registry
add_text(slide, 1.5, 2.0, 5, 0.5, "PublisherRegistry.sol", size=20, bold=True, color=HEDERA_GREEN)
add_bullet_slide(slide, 1.5, 2.7, 5, [
    "registerEndpoint(url, price, paymaster)",
    "On-chain endpoint CRUD + call tracking",
    "Per-endpoint: price, active status, revenue",
    "Publisher ownership verification",
    "Deployed: 0xFBCee3E39A...6189",
], size=13, color=GRAY)

# Right — Paymaster
add_text(slide, 7, 2.0, 5.5, 0.5, "AgentGatePaymaster.sol", size=20, bold=True, color=ACCENT)
add_bullet_slide(slide, 7, 2.7, 5.5, [
    "ERC-4337 v0.7 shared paymaster",
    "Per-endpoint gas budgets (endpointBalance)",
    "Configurable gas share % (0-100%)",
    "fundAndSetGasShare(url, bps) payable",
    "paymasterAndData[52:84] = endpointHash",
    "Deployed: 0xfbC79b8d8b...62a1",
], size=13, color=GRAY)

add_text(slide, 1.5, 5.5, 10, 0.5, "30/30 tests passing (Hardhat + Chai) — all on Hedera Testnet", size=16, bold=True, color=GREEN)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Publisher Dashboard
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 1.5, 0.8, 10, 0.8, "Publisher Dashboard", size=36, bold=True, color=WHITE)
add_text(slide, 1.5, 1.7, 10, 0.5, "React + Vite — connects to MetaMask, reads on-chain state", size=16, color=GRAY)

add_bullet_slide(slide, 1.5, 2.6, 5.5, [
    "📝 Publish Tab",
    "  — Endpoint name, backend URL, auth headers",
    "  — Price per call (USD → HBAR)",
    "  — Gas budget deposit + share %",
    "  — WorldID toggle per endpoint",
    "  — Auto-assigned proxy URL",
    "",
    "⚙️ Manage Tab",
    "  — Call stats (total, free-trial, paid, agents)",
    "  — Gas budget top-up",
    "  — Proxy config (backend URL, headers)",
], size=13, color=GRAY)

add_bullet_slide(slide, 7.5, 2.6, 5, [
    "📊 Dashboard Tab",
    "  — Live chain data (contracts, balances)",
    "  — Registered endpoints with WorldID badges",
    "  — Endpoint names + prices",
    "",
    "⚡ Flow Tab",
    "  — Visual x402 flow diagram",
    "",
    "🔗 All contract links → HashScan",
    "🔗 All tx hashes → verifiable",
], size=13, color=GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Live Demo
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_tag(slide, 10.3, 0.5, "🟣 World", PURPLE)
add_tag(slide, 11.5, 0.5, "🟢 Hedera", HEDERA_GREEN)
add_text(slide, 1.5, 0.8, 10, 0.8, "Live Demo — End to End", size=36, bold=True, color=WHITE)
add_text(slide, 1.5, 1.7, 10, 0.5, "Everything real — HBAR on testnet, OpenAI API, World Chain AgentBook", size=16, color=GREEN)

steps = [
    ("1", "Agent calls /api/proxy/2", "→ HTTP 402 (pay 0.23 HBAR)", GRAY),
    ("2", "Agent signs SIWE challenge", "→ WorldID AgentKit proof", PURPLE),
    ("3", "Agent sends 0.23 HBAR on Hedera", "→ tx confirmed in 3s", HEDERA_GREEN),
    ("4", "Agent retries with payment + agentkit", "→ HTTP 200 ✅", GREEN),
    ("5", "AgentGate forwards to OpenAI", "→ API key injected server-side", ACCENT),
    ("6", "OpenAI responds", "→ \"Hello, how are you?\"", WHITE),
]

for i, (num, title, detail, color) in enumerate(steps):
    y = 2.8 + i * 0.65
    add_text(slide, 1.5, y, 0.5, 0.5, num, size=24, bold=True, color=color)
    add_text(slide, 2.2, y, 5, 0.35, title, size=16, bold=True, color=WHITE)
    add_text(slide, 2.2, y + 0.3, 5, 0.3, detail, size=13, color=GRAY)

# Right side — tx proof
add_text(slide, 7.5, 2.8, 5, 0.4, "Verifiable on HashScan:", size=14, bold=True, color=HEDERA_GREEN)
add_bullet_slide(slide, 7.5, 3.4, 5, [
    "Registry deploy tx: 0x062c3f7d...",
    "Paymaster deploy tx: 0x9b88f4be...",
    "registerEndpoint tx: 0x1801b7a8...",
    "fundAndSetGasShare tx: 0xe7ce811b...",
    "Agent HBAR payment: 0xf7b32626...",
    "",
    "Agent balance: started 99.67 HBAR",
    "After 3 calls: 98.25 HBAR",
    "Real money, real blockchain.",
], size=12, color=GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Architecture
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_tag(slide, 10.3, 0.5, "🟣 World", PURPLE)
add_tag(slide, 11.5, 0.5, "🟢 Hedera", HEDERA_GREEN)
add_text(slide, 1.5, 0.8, 10, 0.8, "Architecture", size=36, bold=True, color=WHITE)

# Architecture as text (clean monospace look)
arch_lines = [
    "┌──────────┐          ┌────────────────────┐          ┌──────────────┐",
    "│          │   402    │                    │  forward  │              │",
    "│  Agent   │ ───────→ │    AgentGate       │ ────────→ │  OpenAI /    │",
    "│          │ ← pay →  │    (Hono + x402)   │ ← resp ← │  Any API     │",
    "└──────────┘          └─────────┬──────────┘          └──────────────┘",
    "                                │",
    "               ┌────────────────┼────────────────┐",
    "               │                │                │",
    "        PublisherRegistry   Paymaster      World Chain",
    "        (Hedera Testnet)   (ERC-4337)    (AgentBook)",
]

for i, line in enumerate(arch_lines):
    add_text(slide, 1.2, 2.2 + i * 0.38, 11, 0.4, line, size=13, color=GRAY, font_name="Courier New")

# Tech stack
add_text(slide, 1.5, 6.2, 10, 0.4, "Stack: Hono · Viem · Hardhat · React · Vite · x402 · @worldcoin/agentkit", size=12, color=DARK_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Why We Win
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 1.5, 0.8, 10, 0.8, "Why AgentGate", size=36, bold=True, color=WHITE)

# Two columns
add_text(slide, 1.5, 2.0, 5, 0.5, "For World Prize", size=22, bold=True, color=PURPLE)
add_bullet_slide(slide, 1.5, 2.7, 5, [
    "✅ AgentKit integrated — not just World ID",
    "✅ SIWE signing + AgentBook verification",
    "✅ Free-trial for human-backed agents",
    "✅ Per-endpoint WorldID toggle (publisher choice)",
    "✅ Distinguishes humans from bots at API level",
    "✅ Real World Chain lookups — zero mocks",
], size=14, color=GRAY)

add_text(slide, 7, 2.0, 5.5, 0.5, "For Hedera Prize", size=22, bold=True, color=HEDERA_GREEN)
add_bullet_slide(slide, 7, 2.7, 5.5, [
    "✅ x402 payment protocol on Hedera",
    "✅ Real HBAR payments (testnet, verifiable)",
    "✅ Mirror Node payment verification",
    "✅ 2 Solidity contracts deployed + 30 tests",
    "✅ ERC-4337 paymaster with gas sponsorship",
    "✅ Pay-per-request API access for AI agents",
], size=14, color=GRAY)

add_text(slide, 1.5, 5.8, 10, 0.6, "AgentGate is the bridge: World verifies who you are. Hedera handles how you pay.", size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Thank You
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 1.5, 2.5, 10, 1.2, "AgentGate", size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, 1.5, 3.7, 10, 0.6, "Pay-per-call API Marketplace for Human-Backed AI Agents", size=20, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(slide, 1.5, 4.8, 10, 0.5, "github.com/sliponit/agentgate", size=14, color=GRAY, align=PP_ALIGN.CENTER)
add_tag(slide, 5.2, 5.8, "🟣 World", PURPLE)
add_tag(slide, 6.5, 5.8, "🟢 Hedera", HEDERA_GREEN)

# Save
output = "/Users/gdonnoh/Desktop/agentgate/AgentGate_Presentation.pptx"
prs.save(output)
print(f"✅ Saved to {output}")
print(f"   {len(prs.slides)} slides")
